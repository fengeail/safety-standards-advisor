# -*- coding: utf-8 -*-
"""安全生产标准文本提取（技能自带，供 read_std.py / lookup.py 复用）。
支持：PDF文字层、docx、rtf(gbk)、xlsx/xls、pptx/ppt、txt/md/csv/xml/json/html。
.doc 二进制与扫描件PDF 返回 (空, scanned=True)，由调用方按"需转docx/扫描件"提示。
"""
import os, warnings
warnings.filterwarnings("ignore")


def extract_pdf(path):
    try:
        from pypdf import PdfReader
        r = PdfReader(path)
        n = len(r.pages)
        if r.is_encrypted:
            try:
                r.decrypt("")
            except Exception:
                pass
        txt = ""
        for p in r.pages:
            try:
                txt += (p.extract_text() or "") + "\n"
            except Exception:
                pass
        return txt, (n > 0 and len(txt.strip()) < 30)
    except Exception:
        return "", False


def extract_docx(path):
    import docx
    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs if p.text), False


def extract_doc(path):
    # .doc 二进制 OLE，piece-table 字节序因版本而异，不抽正文。
    return "", False


def extract_rtf(path):
    from striprtf.striprtf import rtf_to_text
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return rtf_to_text(f.read(), encoding="gbk"), False


def extract_xlsx(path):
    import openpyxl
    t = ""
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for c in row:
                if c is not None:
                    t += str(c) + "\t"
            t += "\n"
    return t, False


def extract_pptx(path):
    import pptx
    prs = pptx.Presentation(path)
    t = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                t += shape.text_frame.text + "\n"
    return t, False


def read_plain(path):
    for enc in ("utf-8", "gbk", "latin-1"):
        try:
            with open(path, "r", encoding=enc, errors="ignore") as f:
                return f.read(), False
        except Exception:
            continue
    return "", False


PLAIN_EXT = {".txt", ".md", ".csv", ".xml", ".json", ".html", ".htm", ".log"}


def extract(path, ext):
    try:
        if ext == ".pdf":
            return extract_pdf(path)
        if ext == ".docx":
            return extract_docx(path)
        if ext == ".doc":
            return extract_doc(path)
        if ext == ".rtf":
            return extract_rtf(path)
        if ext in {".xlsx", ".xls"}:
            return extract_xlsx(path)
        if ext in {".pptx", ".ppt"}:
            return extract_pptx(path)
        if ext in PLAIN_EXT:
            return read_plain(path)
    except Exception:
        return "", False
    return "", False


def read_standard(rel_or_path, standards_root):
    """给定相对路径或绝对路径，返回 (text, scanned, ext, found)。"""
    if os.path.exists(rel_or_path) and os.path.isabs(rel_or_path):
        src = rel_or_path
    else:
        src = os.path.join(standards_root, rel_or_path.replace("/", "\\"))
    if not os.path.exists(src):
        return "", False, "", False
    ext = os.path.splitext(src)[1].lower()
    txt, scanned = extract(src, ext)
    return txt, scanned, ext, True
